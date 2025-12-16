"""
PowerPoint AI Assistant - Fastlane Alternative
Generates and enhances PowerPoint presentations using AI
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import google.generativeai as genai
from dotenv import load_dotenv

# Load API key
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:
    raise ValueError("GEMINI_API_KEY not found in environment variables")

genai.configure(api_key=GEMINI_API_KEY)


class PowerPointAI:
    def __init__(self):
        self.model = genai.GenerativeModel("gemini-2.0-flash-exp")
    
    def generate_presentation_outline(self, topic, num_slides=5):
        """Generate a presentation outline from a topic"""
        prompt = f"""
        Create a professional presentation outline for: "{topic}"
        
        Requirements:
        - {num_slides} slides total
        - Include title slide
        - Each slide should have:
          * Clear title
          * 3-5 bullet points of content
          * Brief speaker notes
        
        Format as JSON:
        {{
          "title": "Presentation Title",
          "slides": [
            {{
              "title": "Slide Title",
              "content": ["Point 1", "Point 2", "Point 3"],
              "notes": "Speaker notes here"
            }}
          ]
        }}
        """
        
        response = self.model.generate_content(prompt)
        import json
        # Parse JSON from response
        text = response.text.strip()
        # Remove markdown code blocks if present
        if text.startswith("```json"):
            text = text[7:]
        if text.startswith("```"):
            text = text[3:]
        if text.endswith("```"):
            text = text[:-3]
        return json.loads(text.strip())
    
    def enhance_slide_content(self, slide_title, bullet_points):
        """Enhance existing slide content with better wording"""
        prompt = f"""
        Improve this slide content to be more professional and impactful:
        
        Title: {slide_title}
        Content:
        {chr(10).join(f"- {point}" for point in bullet_points)}
        
        Return improved content in JSON format:
        {{
          "title": "Improved Title",
          "content": ["Improved point 1", "Improved point 2", ...]
        }}
        
        Keep it concise and professional.
        """
        
        response = self.model.generate_content(prompt)
        import json
        text = response.text.strip()
        if text.startswith("```json"):
            text = text[7:]
        if text.startswith("```"):
            text = text[3:]
        if text.endswith("```"):
            text = text[:-3]
        return json.loads(text.strip())
    
    def create_presentation(self, outline):
        """Create a PowerPoint presentation from outline"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define color scheme
        primary_color = RGBColor(31, 78, 120)  # Dark blue
        accent_color = RGBColor(68, 114, 196)  # Medium blue
        text_color = RGBColor(255, 255, 255)   # White
        bg_color = RGBColor(240, 240, 240)     # Light gray
        
        # Title slide with custom design
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Add background rectangle
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = primary_color
        bg_shape.line.fill.background()
        
        # Add title text box
        title_left = Inches(1)
        title_top = Inches(2.5)
        title_width = Inches(8)
        title_height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = outline["title"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = text_color
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(title_left, Inches(4.2), title_width, Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "AI-Generated Presentation"
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(20)
        subtitle_frame.paragraphs[0].font.color.rgb = text_color
        
        # Content slides with styled design
        for slide_data in outline["slides"]:
            slide = prs.slides.add_slide(blank_layout)
            
            # Add gradient background
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                prs.slide_width, prs.slide_height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = bg_color
            bg_shape.line.fill.background()
            
            # Add header bar
            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                prs.slide_width, Inches(1.2)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = accent_color
            header_shape.line.fill.background()
            
            # Add title in header
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = text_color
            
            # Add content box with rounded corners
            content_left = Inches(0.8)
            content_top = Inches(2)
            content_width = Inches(8.4)
            content_height = Inches(4.5)
            
            content_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                content_left, content_top,
                content_width, content_height
            )
            content_shape.fill.solid()
            content_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            content_shape.line.color.rgb = accent_color
            content_shape.line.width = Pt(2)
            
            # Add content text
            content_box = slide.shapes.add_textbox(
                content_left + Inches(0.3),
                content_top + Inches(0.3),
                content_width - Inches(0.6),
                content_height - Inches(0.6)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, point in enumerate(slide_data["content"]):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f"• {point}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.space_after = Pt(12)
                p.line_spacing = 1.2
            
            # Add speaker notes if present
            if "notes" in slide_data and slide_data["notes"]:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide_data["notes"]
        
        return prs
    
    def enhance_existing_presentation(self, input_path, output_path):
        """Enhance an existing PowerPoint presentation"""
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            # Skip title slides
            if not slide.shapes.title:
                continue
            
            title_text = slide.shapes.title.text
            
            # Find text content
            bullet_points = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text and paragraph.text != title_text:
                            bullet_points.append(paragraph.text)
            
            if bullet_points:
                # Enhance content
                enhanced = self.enhance_slide_content(title_text, bullet_points)
                
                # Update slide
                slide.shapes.title.text = enhanced["title"]
                
                # Update content
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        for point in enhanced["content"]:
                            p = text_frame.add_paragraph()
                            p.text = point
                            p.level = 0
                            p.font.size = Pt(18)
                        break
        
        prs.save(output_path)
        return output_path


def main():
    """Command-line interface"""
    import argparse
    
    parser = argparse.ArgumentParser(description="PowerPoint AI Assistant")
    subparsers = parser.add_subparsers(dest="command", help="Available commands")
    
    # Generate command
    generate_parser = subparsers.add_parser("generate", help="Generate new presentation")
    generate_parser.add_argument("topic", help="Presentation topic")
    generate_parser.add_argument("-n", "--num-slides", type=int, default=5, help="Number of slides")
    generate_parser.add_argument("-o", "--output", default="presentation.pptx", help="Output file")
    
    # Enhance command
    enhance_parser = subparsers.add_parser("enhance", help="Enhance existing presentation")
    enhance_parser.add_argument("input", help="Input PowerPoint file")
    enhance_parser.add_argument("-o", "--output", help="Output file (default: input_enhanced.pptx)")
    
    args = parser.parse_args()
    
    ai = PowerPointAI()
    
    if args.command == "generate":
        print(f"Generating presentation on: {args.topic}")
        outline = ai.generate_presentation_outline(args.topic, args.num_slides)
        prs = ai.create_presentation(outline)
        prs.save(args.output)
        print(f"✓ Presentation saved to: {args.output}")
    
    elif args.command == "enhance":
        output = args.output or args.input.replace(".pptx", "_enhanced.pptx")
        print(f"Enhancing presentation: {args.input}")
        ai.enhance_existing_presentation(args.input, output)
        print(f"✓ Enhanced presentation saved to: {output}")
    
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
